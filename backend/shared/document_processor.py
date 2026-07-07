"""
Document processing module - extracts text from PDF, Office, plain text,
legacy Office, CSV/HTML, email, and other common document formats.
"""
import os
import hashlib
from pathlib import Path
from typing import List, Tuple
import config

# Document processing imports
try:
    import PyPDF2
except ImportError:
    PyPDF2 = None

try:
    from docx import Document as DocxDocument
except ImportError:
    DocxDocument = None

try:
    from pptx import Presentation
except ImportError:
    Presentation = None

try:
    import fitz  # PyMuPDF
except ImportError:
    fitz = None

try:
    from PIL import Image as PILImage
except ImportError:
    PILImage = None

try:
    import sharepoint2text
except ImportError:
    sharepoint2text = None

# File types routed to sharepoint2text (see config._GENERIC_EXTENSIONS).
_GENERIC_EXTRACTOR_TYPES = {ext.lstrip(".") for ext in config._GENERIC_EXTENSIONS}

# Vector metafile formats (common for charts pasted from old Office versions)
# that PIL/OpenAI's vision API can't read directly as raster images - skipped
# rather than sent as broken images.
_UNSUPPORTED_VISION_EXTS = {"emf", "wmf"}


def extract_text_from_pdf(file_path: str) -> str:
    """Extract text from PDF file"""
    if PyPDF2 is None:
        raise ImportError("PyPDF2 is required for PDF processing. Install with: pip install PyPDF2")
    
    text = ""
    try:
        with open(file_path, 'rb') as file:
            pdf_reader = PyPDF2.PdfReader(file)
            for page in pdf_reader.pages:
                text += page.extract_text() + "\n"
    except Exception as e:
        raise Exception(f"Error reading PDF: {str(e)}")
    
    return text.strip()


def _limit_reached(count: int, max_images: int) -> bool:
    """max_images <= 0 means unlimited (see config.VISION_MAX_IMAGES_PER_DOC)."""
    return max_images > 0 and count >= max_images


def _raster_dimensions(image_bytes: bytes):
    """Return (width, height) for a raster image via Pillow, or None if it
    can't be read as one (e.g. a legacy WMF/EMF vector metafile chart)."""
    if PILImage is None or not image_bytes:
        return None
    try:
        import io
        with PILImage.open(io.BytesIO(image_bytes)) as img:
            return img.size
    except Exception:
        return None


def extract_images_from_pdf(
    file_path: str,
    min_dim: int = None,
    max_images: int = None,
) -> List[Tuple[str, bytes, str]]:
    """Pull embedded images out of a PDF for optional vision description.

    Returns a list of (location_label, image_bytes, image_ext) tuples, e.g.
    ("page 3", b"...", "png"). Filters out tiny images (icons/logos/dividers)
    and exact duplicates (the same image repeated across pages/slide
    masters), and stops once max_images is reached, to bound cost/time on
    image-heavy documents.
    """
    if fitz is None:
        return []

    min_dim = config.VISION_MIN_IMAGE_DIM if min_dim is None else min_dim
    max_images = config.VISION_MAX_IMAGES_PER_DOC if max_images is None else max_images

    results: List[Tuple[str, bytes, str]] = []
    seen_hashes = set()

    try:
        doc = fitz.open(file_path)
        try:
            for page_index in range(len(doc)):
                if _limit_reached(len(results), max_images):
                    break
                page = doc[page_index]
                for img in page.get_images(full=True):
                    if _limit_reached(len(results), max_images):
                        break
                    xref = img[0]
                    try:
                        base_image = doc.extract_image(xref)
                    except Exception:
                        continue
                    if not base_image:
                        continue

                    ext = base_image.get("ext", "png")
                    if ext.lower() in _UNSUPPORTED_VISION_EXTS:
                        continue

                    width = base_image.get("width", 0)
                    height = base_image.get("height", 0)
                    if width < min_dim or height < min_dim:
                        continue

                    image_bytes = base_image.get("image")
                    if not image_bytes:
                        continue

                    digest = hashlib.sha1(image_bytes).hexdigest()
                    if digest in seen_hashes:
                        continue
                    seen_hashes.add(digest)

                    results.append((f"page {page_index + 1}", image_bytes, ext))
        finally:
            doc.close()
    except Exception as e:
        print(f"[VISION] Failed to extract images from {file_path}: {e}")

    return results


def extract_images_from_pptx(
    file_path: str,
    min_dim: int = None,
    max_images: int = None,
) -> List[Tuple[str, bytes, str]]:
    """Pull picture images out of a PPTX for optional vision description.

    Returns (location_label, image_bytes, image_ext) tuples, e.g.
    ("slide 4", b"...", "png"). Note: native PowerPoint charts/SmartArt (not
    pasted-in pictures) aren't stored as images at all and won't be caught
    here - only actual picture shapes are.
    """
    if Presentation is None:
        return []

    min_dim = config.VISION_MIN_IMAGE_DIM if min_dim is None else min_dim
    max_images = config.VISION_MAX_IMAGES_PER_DOC if max_images is None else max_images

    results: List[Tuple[str, bytes, str]] = []
    seen_hashes = set()

    try:
        prs = Presentation(file_path)
        for slide_index, slide in enumerate(prs.slides, start=1):
            if _limit_reached(len(results), max_images):
                break
            for shape in slide.shapes:
                if _limit_reached(len(results), max_images):
                    break
                image = getattr(shape, "image", None)
                if image is None:
                    continue
                try:
                    image_bytes = image.blob
                    ext = (image.ext or "png").lower()
                except Exception:
                    continue

                if ext in _UNSUPPORTED_VISION_EXTS:
                    continue

                digest = hashlib.sha1(image_bytes).hexdigest()
                if digest in seen_hashes:
                    continue

                dims = _raster_dimensions(image_bytes)
                if dims is None:
                    continue
                width, height = dims
                if width < min_dim or height < min_dim:
                    continue

                seen_hashes.add(digest)
                results.append((f"slide {slide_index}", image_bytes, ext))
    except Exception as e:
        print(f"[VISION] Failed to extract images from {file_path}: {e}")

    return results


def extract_images_from_docx(
    file_path: str,
    min_dim: int = None,
    max_images: int = None,
) -> List[Tuple[str, bytes, str]]:
    """Pull embedded images out of a DOCX for optional vision description.

    Returns (location_label, image_bytes, image_ext) tuples. Word doesn't
    store page breaks as fixed data (pagination is a rendering-time concept),
    so images are labeled by order rather than a real page number.
    """
    if DocxDocument is None:
        return []

    min_dim = config.VISION_MIN_IMAGE_DIM if min_dim is None else min_dim
    max_images = config.VISION_MAX_IMAGES_PER_DOC if max_images is None else max_images

    results: List[Tuple[str, bytes, str]] = []
    seen_hashes = set()

    try:
        doc = DocxDocument(file_path)
        order = 0
        for rel in doc.part.rels.values():
            if _limit_reached(len(results), max_images):
                break
            if "image" not in rel.reltype:
                continue
            try:
                image_bytes = rel.target_part.blob
                content_type = rel.target_part.content_type  # e.g. "image/png"
            except Exception:
                continue

            ext = content_type.split("/")[-1].lower() if content_type else "png"
            if ext in _UNSUPPORTED_VISION_EXTS:
                continue

            digest = hashlib.sha1(image_bytes).hexdigest()
            if digest in seen_hashes:
                continue

            dims = _raster_dimensions(image_bytes)
            if dims is None:
                continue
            width, height = dims
            if width < min_dim or height < min_dim:
                continue

            seen_hashes.add(digest)
            order += 1
            results.append((f"image {order}", image_bytes, ext))
    except Exception as e:
        print(f"[VISION] Failed to extract images from {file_path}: {e}")

    return results


def extract_images_from_document(
    file_path: str,
    file_type: str,
    min_dim: int = None,
    max_images: int = None,
) -> List[Tuple[str, bytes, str]]:
    """Dispatch to the right image extractor based on file type. Returns
    (location_label, image_bytes, image_ext) tuples, or [] for unsupported
    types (txt has no images; anything else not implemented yet).
    """
    file_type = file_type.lower().lstrip('.')
    if file_type == "pdf":
        return extract_images_from_pdf(file_path, min_dim, max_images)
    elif file_type == "pptx":
        return extract_images_from_pptx(file_path, min_dim, max_images)
    elif file_type == "docx":
        return extract_images_from_docx(file_path, min_dim, max_images)
    return []


def extract_text_from_docx(file_path: str) -> str:
    """Extract text from DOCX file"""
    if DocxDocument is None:
        raise ImportError("python-docx is required for DOCX processing. Install with: pip install python-docx")
    
    try:
        doc = DocxDocument(file_path)
        text = "\n".join([paragraph.text for paragraph in doc.paragraphs])
    except Exception as e:
        raise Exception(f"Error reading DOCX: {str(e)}")
    
    return text.strip()


def extract_text_from_pptx(file_path: str) -> str:
    """Extract text from PPTX file"""
    if Presentation is None:
        raise ImportError("python-pptx is required for PPTX processing. Install with: pip install python-pptx")
    
    try:
        prs = Presentation(file_path)
        text_parts = []
        for slide in prs.slides:
            for shape in slide.shapes:
                if hasattr(shape, "text"):
                    text_parts.append(shape.text)
        text = "\n".join(text_parts)
    except Exception as e:
        raise Exception(f"Error reading PPTX: {str(e)}")
    
    return text.strip()


def extract_text_from_txt(file_path: str) -> str:
    """Extract text from TXT file"""
    try:
        with open(file_path, 'r', encoding='utf-8', errors='ignore') as file:
            text = file.read()
    except Exception as e:
        raise Exception(f"Error reading TXT: {str(e)}")
    
    return text.strip()


def extract_text_generic(file_path: str) -> str:
    """Extract text from legacy Office (.doc/.ppt/.xls/.rtf), CSV, HTML, and
    OpenOffice (.odt/.odp/.ods) files via sharepoint2text - pure Python, no
    LibreOffice or other system dependency required."""
    if sharepoint2text is None:
        raise ImportError(
            "sharepoint-to-text is required for this file type. "
            "Install with: pip install sharepoint-to-text"
        )
    try:
        result = next(sharepoint2text.read_file(file_path, ignore_images=True))
        text = result.get_full_text()
    except StopIteration:
        text = ""
    except Exception as e:
        raise Exception(f"Error reading file: {str(e)}")

    return (text or "").strip()


def extract_text(file_path: str, file_type: str) -> str:
    """Extract text from any supported file type"""
    file_type = file_type.lower().lstrip('.')
    
    if file_type == 'pdf':
        return extract_text_from_pdf(file_path)
    elif file_type == 'docx':
        return extract_text_from_docx(file_path)
    elif file_type == 'pptx':
        return extract_text_from_pptx(file_path)
    elif file_type in ('txt', 'md'):
        return extract_text_from_txt(file_path)
    elif file_type in _GENERIC_EXTRACTOR_TYPES:
        return extract_text_generic(file_path)
    else:
        raise ValueError(f"Unsupported file type: {file_type}")


def split_text_into_chunks(text: str, chunk_size: int = None, chunk_overlap: int = None) -> List[str]:
    """
    Split text into chunks for RAG processing
    
    Args:
        text: The text to split
        chunk_size: Size of each chunk in characters (default from config)
        chunk_overlap: Overlap between chunks in characters (default from config)
    
    Returns:
        List of text chunks
    """
    if chunk_size is None:
        chunk_size = config.CHUNK_SIZE
    if chunk_overlap is None:
        chunk_overlap = config.CHUNK_OVERLAP
    
    if len(text) <= chunk_size:
        return [text]
    
    # Ensure overlap is less than chunk size
    if chunk_overlap >= chunk_size:
        chunk_overlap = chunk_size // 4  # Default to 25% overlap
    
    chunks = []
    start = 0
    max_iterations = len(text) // (chunk_size - chunk_overlap) + 10  # Safety limit
    iteration = 0
    
    while start < len(text) and iteration < max_iterations:
        iteration += 1
        end = min(start + chunk_size, len(text))
        
        # Try to break at sentence boundary (only if not at end of text)
        if end < len(text):
            # Look for sentence endings
            for punct in ['. ', '.\n', '! ', '!\n', '? ', '?\n']:
                last_punct = text.rfind(punct, start, end)
                if last_punct != -1:
                    end = last_punct + 2  # Include punctuation and space
                    break
        
        # Ensure we don't go past text length
        end = min(end, len(text))
        
        # Get chunk (don't strip yet - we need to preserve position)
        chunk = text[start:end]
        
        # Only add non-empty chunks
        if chunk.strip():
            chunks.append(chunk.strip())
        
        # If we've reached the end of text, we're done
        if end >= len(text):
            break
        
        # Move start position with overlap
        # Ensure we make progress (at least chunk_size - overlap characters)
        new_start = end - chunk_overlap
        if new_start <= start:
            # Force progress - move forward by at least 1/4 of chunk size
            new_start = start + max(1, (chunk_size - chunk_overlap) // 4)
        
        # Don't go past end of text
        if new_start >= len(text):
            break
            
        start = new_start
    
    return chunks

